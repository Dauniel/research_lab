from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from datetime import date
import copy

# ── palette ──────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0xFF, 0xFF, 0xFF)   # white background
ACCENT    = RGBColor(0x1A, 0x73, 0xE8)   # blue
ACCENT2   = RGBColor(0x18, 0x96, 0x5C)   # green
WHITE     = RGBColor(0x1A, 0x1A, 0x1A)   # near-black text
LIGHT_GRAY= RGBColor(0x55, 0x55, 0x55)   # dark gray
YELLOW    = RGBColor(0xD6, 0x74, 0x00)   # amber
RED_SOFT  = RGBColor(0xC6, 0x28, 0x28)   # red

W, H = Inches(13.33), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(blank_layout)
    # dark background rectangle
    bg = sl.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return sl

def txb(sl, text, x, y, w, h, size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    box = sl.shapes.add_textbox(x, y, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box

def rect(sl, x, y, w, h, fill, line_color=None):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
    else:
        sh.line.fill.background()
    return sh

def accent_bar(sl, color=ACCENT):
    rect(sl, 0, Inches(6.9), W, Inches(0.08), color)

def bullet_box(sl, items, x, y, w, h, size=20, color=WHITE, spacing=Pt(6)):
    box = sl.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = color

def divider(sl, y, color=ACCENT):
    rect(sl, Inches(0.5), y, Inches(12.33), Inches(0.03), color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()

# top color band
rect(sl, 0, 0, W, Inches(0.5), ACCENT)

txb(sl, "Automated Condensate Partition Coefficient Pipeline",
    Inches(0.6), Inches(1.2), Inches(12), Inches(1.8),
    size=38, bold=True, color=WHITE)

txb(sl, "Spring 2026 Progress — Week 7",
    Inches(0.6), Inches(3.0), Inches(8), Inches(0.6),
    size=24, color=ACCENT)

txb(sl, "Daniel Chang  ·  Franco Lab  ·  C&S BIO 199/197",
    Inches(0.6), Inches(3.8), Inches(10), Inches(0.5),
    size=18, color=LIGHT_GRAY)

txb(sl, "May 7, 2026",
    Inches(0.6), Inches(4.3), Inches(6), Inches(0.4),
    size=16, color=LIGHT_GRAY)

accent_bar(sl, ACCENT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda / overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Since Last Meeting (April 22)", Inches(0.5), Inches(0.25),
    Inches(12), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

items = [
    "1.  Confirmed reference PC uses background subtraction  →  added to pipeline",
    "2.  Built spring_implementation/pipeline.py  (full Cellpose 3, 3D end-to-end)",
    "3.  Fixed nuclei over-segmentation with 3D connected-component relabeling",
    "4.  Solved dilute density instability  →  lowest-50-patch selection",
    "5.  Single-ROI result: PC = 6.297  vs  reference 6.32  (−0.4% error)",
    "6.  Batch validation across 30 JABr cells  →  r = 0.735,  RMSE = 2.815,  bias = +2.8%",
    "7.  Discovered Cellpose over-draws boundaries  →  top-75% voxel trim as fix",
]
bullet_box(sl, items, Inches(0.6), Inches(1.1), Inches(12), Inches(5.5), size=21)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Pipeline overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Pipeline Overview (pipeline.py)", Inches(0.5), Inches(0.25),
    Inches(12), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

steps = [
    ("1  Load",      "Condensate (Ch2) + Nuclei (Ch1) TIF stacks  |  55 × 185 × 259 px"),
    ("2  Denoise",   "Cellpose 3 DenoiseModel on every Z-slice  →  sharper boundaries\n    (raw pixels unchanged — denoising only guides masks)"),
    ("3  Segment",   "Cellpose 3 cyto3, do_3D=True  →  full 3D tracking\n    Nuclei post-proc: binary mask → 3D connected components → drop noise\n    76 Cellpose labels  →  5 clean nuclei"),
    ("4  Measure",   "regionprops_table: area, centroid, mean intensity per object per Z"),
    ("5  3D Volume", "Count voxels per label across all 55 slices"),
    ("6  PC",        "Background-subtracted Fabrini formula\n    cond_density = mean top-75% bright voxels in cond ∩ nuc mask\n    dil_density  = mean of 50 lowest-intensity 10×10×10 patches in dilute region\n    PC = cond_density / dil_density"),
]

y = Inches(1.1)
for label, desc in steps:
    rect(sl, Inches(0.5), y, Inches(1.5), Inches(0.55), ACCENT2)
    txb(sl, label, Inches(0.55), y + Inches(0.05), Inches(1.4), Inches(0.5),
        size=15, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    txb(sl, desc, Inches(2.2), y, Inches(10.8), Inches(0.65),
        size=15, color=WHITE)
    y += Inches(0.95)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Single-ROI result
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Single-ROI Validation  (Sample2_5_1)", Inches(0.5), Inches(0.25),
    Inches(12), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

# big number
txb(sl, "PC = 6.297", Inches(1.0), Inches(1.3), Inches(6), Inches(1.2),
    size=60, bold=True, color=ACCENT)
txb(sl, "Reference = 6.32  (−0.4% error)",
    Inches(1.0), Inches(2.5), Inches(8), Inches(0.6),
    size=26, color=YELLOW)

divider(sl, Inches(3.3), ACCENT2)

txb(sl, "Three improvements that closed the gap (3.3 → 6.3):",
    Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
    size=20, bold=True, color=LIGHT_GRAY)

items = [
    "Background subtraction  (April 25) — reference was bg-subtracted; pipeline wasn't.  Largest single jump.",
    "Connected-component nuclei relabeling  (April 30) — 76 Cellpose fragments → 5 true nuclei.  Fixes dilute-phase pixel coverage.",
    "Lowest-50-patch dilute density  (April 30) — replaces unstable single random patch (PC swing: 4.3–6.0 across seeds).",
]
bullet_box(sl, items, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.5), size=19)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Top-75% fix
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Key Finding: Cellpose Over-draws Boundaries", Inches(0.5), Inches(0.25),
    Inches(12), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

txb(sl, "Diagnosis (Sample3_3_11, ref PC = 16.95):",
    Inches(0.5), Inches(1.05), Inches(12), Inches(0.45), size=21, bold=True, color=LIGHT_GRAY)

problem = [
    "Pipeline cond density:   466   →  pipeline PC = 8.32",
    "Reference cond density: 1362   →  reference PC = 16.95",
    "",
    "Cellpose 3D masks include bright core + large dim halo + dark interior pixels.",
    "Manual Imaris boundaries capture only the bright core.",
]
bullet_box(sl, problem, Inches(0.7), Inches(1.55), Inches(7.5), Inches(2.2), size=19)

# table on right
headers = ["Variant", "r", "RMSE", "Bias"]
rows = [
    ["top-10%",  "0.912", "13.10", "+120%"],
    ["top-25%",  "0.865",  "8.17",  "+75%"],
    ["top-50%",  "0.783",  "3.93",  "+30%"],
    ["top-75% ✓","0.735",  "2.82",   "+3%"],
    ["full mean","0.716",  "3.25",  "−16%"],
]

col_w = [Inches(2.0), Inches(0.9), Inches(0.9), Inches(0.9)]
col_x = [Inches(8.4), Inches(10.4), Inches(11.3), Inches(12.2)]
row_h = Inches(0.46)
t_y   = Inches(1.05)

for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    rect(sl, cx, t_y, cw, row_h, ACCENT2)
    txb(sl, hdr, cx + Inches(0.05), t_y + Inches(0.08), cw, row_h,
        size=15, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    fill = RGBColor(0xD9, 0xEA, 0xD3) if row[0].startswith("top-75") else RGBColor(0xF1, 0xF3, 0xF4)
    for ci, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        rect(sl, cx, t_y + row_h * (ri + 1), cw, row_h, fill)
        txb(sl, val, cx + Inches(0.05), t_y + row_h * (ri + 1) + Inches(0.08),
            cw, row_h, size=15, color=RGBColor(0x1A,0x1A,0x1A), align=PP_ALIGN.CENTER)

txb(sl, "top-75% trims the ~25% 'fluff' Cellpose adds beyond a manual boundary.\nNear-zero bias and best RMSE across 29 cells.",
    Inches(0.6), Inches(5.9), Inches(12), Inches(0.9), size=18, color=YELLOW)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Batch results
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Batch Validation — 30 JABr Cells", Inches(0.5), Inches(0.25),
    Inches(12), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

# summary boxes
for i, (label, val, color) in enumerate([
    ("Pearson r", "0.735", ACCENT),
    ("RMSE", "2.815", ACCENT2),
    ("Mean bias", "+2.8%", YELLOW),
    ("Cells processed", "29 / 30", LIGHT_GRAY),
]):
    bx = Inches(0.5) + i * Inches(3.2)
    rect(sl, bx, Inches(1.1), Inches(3.0), Inches(1.3), RGBColor(0xF1, 0xF3, 0xF4))
    txb(sl, val,   bx + Inches(0.1), Inches(1.15), Inches(2.8), Inches(0.75),
        size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(sl, label, bx + Inches(0.1), Inches(1.85), Inches(2.8), Inches(0.4),
        size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txb(sl, "Well-calibrated cells (error < 5%)",
    Inches(0.5), Inches(2.65), Inches(6), Inches(0.4),
    size=18, bold=True, color=LIGHT_GRAY)
good = [
    "Sample1_1_3  ref=19.77  pipe=19.90  (+0.7%)",
    "Sample3_3_3   ref=8.82   pipe=8.74   (−1.0%)",
    "Sample3_3_7   ref=8.68   pipe=8.55   (−1.5%)",
    "Sample1_4_1   ref=6.13   pipe=6.20   (+1.1%)",
    "Sample3_3_9   ref=5.92   pipe=6.13   (+3.6%)",
]
bullet_box(sl, good, Inches(0.6), Inches(3.1), Inches(6.2), Inches(2.5), size=17)

txb(sl, "Known failures (segmentation, not calibration)",
    Inches(7.0), Inches(2.65), Inches(6), Inches(0.4),
    size=18, bold=True, color=RED_SOFT)
bad = [
    "Sample1_1_1   +146%  wrong cell selected",
    "Sample3_3_2    −51%  missed nucleus",
    "Sample3_3_11   −38%  bright cores undersampled",
    "Sample2_5_5    −38%  large FOV, nucleus not detected",
    "Sample3_3_15    NaN  zero condensate in central nuc",
]
bullet_box(sl, bad, Inches(7.1), Inches(3.1), Inches(6.0), Inches(2.5), size=17, color=RED_SOFT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Per-cell comparison table
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Per-Cell Results: Reference vs Pipeline (top-75%)", Inches(0.5), Inches(0.25),
    Inches(12.5), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

# table data: (sample, ref_pc, pipe_pc, error, is_failure)
table_data = [
    ("Sample1_1_3",  "19.77", "19.90",  "+0.7%",  False),
    ("Sample1_4_3",  "13.96", "13.04",  "−6.6%",  False),
    ("Sample3_3_3",   "8.82",  "8.74",  "−1.0%",  False),
    ("Sample3_3_7",   "8.68",  "8.55",  "−1.5%",  False),
    ("Sample3_3_1",   "8.18",  "8.31",  "+1.5%",  False),
    ("Sample3_3_9",   "5.92",  "6.13",  "+3.6%",  False),
    ("Sample1_4_1",   "6.13",  "6.20",  "+1.1%",  False),
    ("Sample3_3_11", "16.95", "10.57", "−38.0%",  True),
    ("Sample3_3_2",   "9.05",  "4.47", "−50.6%",  True),
    ("Sample2_5_5",   "6.73",  "4.20", "−37.6%",  True),
    ("Sample1_1_1",   "4.78", "11.77", "+146%",   True),
    ("Sample3_3_15",  "6.41",   "NaN",    "NaN",  True),
]

headers    = ["Sample", "Reference PC", "Pipeline PC (top-75%)", "Error"]
col_widths = [Inches(3.2), Inches(2.4), Inches(3.6), Inches(2.0)]
col_xs     = [Inches(0.5), Inches(3.7), Inches(6.1), Inches(9.7)]
row_h      = Inches(0.43)
t_y        = Inches(1.05)

HDR_FILL  = ACCENT
HDR_TEXT  = RGBColor(0xFF, 0xFF, 0xFF)
GOOD_FILL = RGBColor(0xE8, 0xF5, 0xE9)
FAIL_FILL = RGBColor(0xFF, 0xEB, 0xEE)
ALT_FILL  = RGBColor(0xF8, 0xF9, 0xFA)

# header row
for hdr, cx, cw in zip(headers, col_xs, col_widths):
    rect(sl, cx, t_y, cw, row_h, HDR_FILL)
    txb(sl, hdr, cx + Inches(0.08), t_y + Inches(0.09), cw, row_h,
        size=15, bold=True, color=HDR_TEXT, align=PP_ALIGN.CENTER)

# data rows
for ri, (sample, ref, pipe, err, fail) in enumerate(table_data):
    fill = FAIL_FILL if fail else (GOOD_FILL if ri % 2 == 0 else ALT_FILL)
    text_color = RED_SOFT if fail else RGBColor(0x1A, 0x1A, 0x1A)
    for val, cx, cw in zip([sample, ref, pipe, err], col_xs, col_widths):
        rect(sl, cx, t_y + row_h * (ri + 1), cw, row_h, fill,
             line_color=RGBColor(0xDD, 0xDD, 0xDD))
        txb(sl, val, cx + Inches(0.08), t_y + row_h * (ri + 1) + Inches(0.08),
            cw, row_h, size=14, color=text_color, align=PP_ALIGN.CENTER)

# legend
rect(sl, Inches(0.5), Inches(6.55), Inches(0.25), Inches(0.2), GOOD_FILL,
     line_color=RGBColor(0xDD,0xDD,0xDD))
txb(sl, "Well-calibrated", Inches(0.8), Inches(6.53), Inches(2.5), Inches(0.3),
    size=13, color=LIGHT_GRAY)
rect(sl, Inches(3.5), Inches(6.55), Inches(0.25), Inches(0.2), FAIL_FILL,
     line_color=RGBColor(0xDD,0xDD,0xDD))
txb(sl, "Segmentation failure", Inches(3.8), Inches(6.53), Inches(3.0), Inches(0.3),
    size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Outstanding issues
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Outstanding Issues", Inches(0.5), Inches(0.25),
    Inches(12), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

issues = [
    ("Nuclei over-segmentation",
     "Cellpose still labels each nucleus as ~15–25 fragments internally.\n"
     "Connected-component fix makes the binary mask correct, but true fix\n"
     "would need 2D+stitch or a custom-trained model."),
    ("top-75% not universally optimal",
     "Calibrated for JABr. If Cellpose draws tighter masks on other constructs\n"
     "(JwtBr, 10ntABr), trimming 25% may over-correct.\n"
     "Would need per-cell mask quality assessment to tune adaptively."),
    ("Single-cell selection heuristic",
     "Central nucleus works for well-cropped ROIs. Fails on wide-FOV images\n"
     "(Sample2_5_5: 32×417×370) where the target cell isn't centered."),
    ("Sample3_3_15 NaN",
     "Central nucleus has zero condensate overlap. Need fallback selection\n"
     "(e.g. max-overlap nucleus) when central nucleus is empty."),
]

y = Inches(1.1)
for title, body in issues:
    rect(sl, Inches(0.5), y, Inches(0.06), Inches(0.9), YELLOW)
    txb(sl, title, Inches(0.7), y, Inches(12), Inches(0.38),
        size=19, bold=True, color=YELLOW)
    txb(sl, body, Inches(0.7), y + Inches(0.38), Inches(12), Inches(0.6),
        size=16, color=LIGHT_GRAY)
    y += Inches(1.35)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Discussion / next steps
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
txb(sl, "Discussion & Next Steps", Inches(0.5), Inches(0.25),
    Inches(12), Inches(0.6), size=30, bold=True, color=ACCENT)
divider(sl, Inches(0.95))

txb(sl, "Questions to discuss:", Inches(0.5), Inches(1.1),
    Inches(12), Inches(0.45), size=22, bold=True, color=LIGHT_GRAY)

questions = [
    "Is top-75% voxel trimming a defensible method to report in the paper,\n"
    "   or should we pursue tighter segmentation instead?",
    "Run batch_compare on other constructs (JwtBr, 10ntABr) to test generalizability?",
    "Outlier cells (Sample2_5_5, Sample3_3_15) — worth fixing, or exclude with justification?",
    "Poster / report scope — what level of batch validation is enough?",
]
bullet_box(sl, questions, Inches(0.7), Inches(1.65), Inches(12), Inches(2.5), size=19)

divider(sl, Inches(4.35), ACCENT2)

txb(sl, "Planned next steps:", Inches(0.5), Inches(4.5),
    Inches(12), Inches(0.45), size=22, bold=True, color=LIGHT_GRAY)
next_steps = [
    "Run batch_compare on JwtBr / 10ntABr constructs",
    "Fix Sample3_3_15 NaN: add fallback nucleus selection",
    "Add discussion + applicability slides to presentation",
    "Explore 2D+stitch approach for nuclei segmentation (if time permits)",
]
bullet_box(sl, next_steps, Inches(0.7), Inches(5.05), Inches(12), Inches(1.8), size=19, color=ACCENT)


# ── save ──────────────────────────────────────────────────────────────────────
out = f"spring_implementation/slides/meetings/meeting_slides_{date.today()}.pptx"
prs.save(out)
print(f"Saved → {out}")
