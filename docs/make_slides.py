from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BG = RGBColor(0xFF, 0xFF, 0xFF)        # white
ACCENT = RGBColor(0x1A, 0x6E, 0xB5)   # blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0x1A, 0x1A, 0x2E)    # dark text
MUTED = RGBColor(0x66, 0x77, 0x88)
GREEN = RGBColor(0x15, 0x80, 0x3D)
YELLOW = RGBColor(0x92, 0x40, 0x00)
RED = RGBColor(0xB9, 0x1C, 0x1C)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def txbox(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def rule(slide, y, color=ACCENT):
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), y, Inches(12.33), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def header(slide, title, subtitle=None):
    txbox(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.65),
          size=32, bold=True, color=ACCENT)
    if subtitle:
        txbox(slide, subtitle, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
              size=16, color=MUTED)
    rule(slide, Inches(1.25))


# ── Slide 1: Title ────────────────────────────────────────────────────────────
slide = add_slide()
txbox(slide, "Spring 2026 Check-in",
      Inches(0.5), Inches(1.8), Inches(12), Inches(0.8),
      size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(slide, "Segmentation Model Survey — JABr Condensates",
      Inches(0.5), Inches(2.65), Inches(12), Inches(0.55),
      size=22, color=LIGHT, align=PP_ALIGN.CENTER)
txbox(slide, "Daniel Chang  ·  C&S BIO 199/197  ·  PI: Elisa Franco  ·  2026-04-23",
      Inches(0.5), Inches(3.3), Inches(12), Inches(0.4),
      size=15, color=MUTED, align=PP_ALIGN.CENTER)

rule(slide, Inches(4.0))

bullets = [
    "Winter pipeline left a PC gap: automated 2.04 vs. manual reference 6.32",
    "Goal this week: survey alternative segmentation models to close the gap",
    "Dataset: JABr ROI Z-stack  (55 slices × 185 × 259, 2-channel fluorescence)",
]
y = Inches(4.2)
for b in bullets:
    txbox(slide, f"·  {b}", Inches(1.5), y, Inches(10), Inches(0.38), size=17, color=LIGHT)
    y += Inches(0.42)

# ── Slide 2: Model Comparison Table ───────────────────────────────────────────
slide = add_slide()
header(slide, "Model Comparison", "Same PC formula, same dataset — 4 new models benchmarked")

cols = [Inches(0.5), Inches(3.8), Inches(8.2), Inches(10.2)]
col_w = [Inches(3.2), Inches(4.3), Inches(1.9), Inches(2.1)]
rows = [
    ("Model", "Approach", "PC", "Runtime", True, ACCENT),
    ("Old Cellpose (2D)", "cyto2, slice-by-slice (baseline)", "2.04", "—", False, LIGHT),
    ("Cellpose 3  ★", "cyto3 + denoising + 3D mode", "3.33", "11 s", False, GREEN),
    ("StarDist 2D", "versatile_fluo, slice-by-slice", "1.91", "6 s", False, LIGHT),
    ("U-FISH †", "ONNX spot detection", "3.81", "2 s", False, YELLOW),
    ("Nellie", "Frangi filter + graph analysis", "3.57", "14 s", False, LIGHT),
    ("Reference", "Manual / ImageJ benchmark", "6.32", "—", False, ACCENT),
]

row_h = Inches(0.52)
y0 = Inches(1.45)
for i, (c0, c1, c2, c3, is_hdr, color) in enumerate(rows):
    y = y0 + i * row_h
    if is_hdr:
        bg_rect = slide.shapes.add_shape(1, Inches(0.4), y, Inches(12.5), row_h)
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = RGBColor(0x1A, 0x6E, 0xB5)
        bg_rect.line.fill.background()
    elif i % 2 == 0:
        bg_rect = slide.shapes.add_shape(1, Inches(0.4), y, Inches(12.5), row_h)
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        bg_rect.line.fill.background()
    for col_x, cw, text in zip(cols, col_w, [c0, c1, c2, c3]):
        txbox(slide, text, col_x, y + Inches(0.06), cw, Inches(0.4),
              size=15 if not is_hdr else 16, bold=is_hdr, color=color)

txbox(slide, "★ Recommended    † U-FISH nuclei mask = Otsu threshold, not a trained model — PC likely inflated",
      Inches(0.5), Inches(6.9), Inches(12), Inches(0.4), size=12, color=MUTED)

# ── Slide 3: Why Cellpose 3 ────────────────────────────────────────────────────
slide = add_slide()
header(slide, "Why Cellpose 3?", "Three concrete changes over the Winter baseline (PC: 2.04 → 3.33)")

items = [
    ("1  Denoising", "DenoiseModel cyto3 runs on every slice before segmentation — suppresses shot noise, sharpens condensate boundaries.\nImportant: PC intensities always computed from raw pixels. Denoising only guides boundary placement."),
    ("2  3D mode", "do_3D=True processes XY, XZ, and YZ planes simultaneously and merges gradient flows before drawing boundaries.\nOld pipeline: 55 independent 2D slices → inconsistent boundaries per condensate → hurts intensity ratio."),
    ("3  cyto3 vs cyto2", "cyto3 trained on a larger, more diverse dataset. Better at small, dim, and irregularly shaped objects.\n3D mode is likely the single largest contributor to the PC improvement."),
]

y = Inches(1.45)
for title, body in items:
    # accent bar
    bar = slide.shapes.add_shape(1, Inches(0.5), y + Inches(0.05), Inches(0.08), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    txbox(slide, title, Inches(0.72), y, Inches(11.8), Inches(0.38), size=18, bold=True, color=ACCENT)
    txbox(slide, body, Inches(0.72), y + Inches(0.35), Inches(11.8), Inches(0.65), size=14, color=LIGHT)
    y += Inches(1.15)

# bottom callout
box = slide.shapes.add_shape(1, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.8))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
box.line.color.rgb = ACCENT
txbox(slide, "All models still fall short of reference PC = 6.32 — gap is due to condensate segmentation sensitivity. "
      "Small boundary errors have an outsized effect on intensity ratios.",
      Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.65), size=14, color=LIGHT)

# ── Slide 4: Comparison Figure ────────────────────────────────────────────────
slide = add_slide()
header(slide, "Segmentation Output — Side-by-Side", "Mid-Z slice, all models")

fig_path = "segmentation_test/outputs/comparison/comparison_figure.png"
try:
    slide.shapes.add_picture(fig_path, Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.8))
except Exception:
    txbox(slide, f"[figure: {fig_path}]", Inches(0.5), Inches(3.5), Inches(12), Inches(0.5),
          size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ── Slide 5: Caveats ──────────────────────────────────────────────────────────
slide = add_slide()
header(slide, "Caveats — Apples-to-Apples Problem", "PC values are not fully comparable across models")

rows2 = [
    ("Model", "Condensate mask", "Nuclei mask", True),
    ("Old Cellpose", "Cellpose cyto2 (2D)", "Cellpose cyto2 (2D)", False),
    ("Cellpose 3", "Cellpose cyto3 (3D)", "Cellpose cyto3 (3D)", False),
    ("StarDist 2D", "StarDist2D", "StarDist2D", False),
    ("U-FISH  ⚠", "U-FISH spot + disk", "Otsu threshold  ← different!", False),
    ("Nellie", "Frangi + graph", "Frangi + graph", False),
]

cols2 = [Inches(0.5), Inches(4.3), Inches(8.5)]
col_w2 = [Inches(3.7), Inches(4.1), Inches(4.3)]
row_h2 = Inches(0.48)
y0 = Inches(1.45)
for i, (c0, c1, c2, is_hdr) in enumerate(rows2):
    y = y0 + i * row_h2
    color = ACCENT if is_hdr else (YELLOW if "⚠" in c0 else LIGHT)
    if is_hdr:
        bg = slide.shapes.add_shape(1, Inches(0.4), y, Inches(12.5), row_h2)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1A, 0x6E, 0xB5)
        bg.line.fill.background()
    elif i % 2 == 0:
        bg = slide.shapes.add_shape(1, Inches(0.4), y, Inches(12.5), row_h2)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        bg.line.fill.background()
    for cx, cw, text in zip(cols2, col_w2, [c0, c1, c2]):
        txbox(slide, text, cx, y + Inches(0.05), cw, Inches(0.38),
              size=15 if not is_hdr else 16, bold=is_hdr, color=color)

txbox(slide, "Fix: use the same Cellpose 3 nuclear mask for all models and only swap the condensate mask — "
      "isolates condensate segmentation as the single variable.",
      Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.55), size=15, bold=True, color=GREEN)

# ── Slide 6: Next Steps ───────────────────────────────────────────────────────
slide = add_slide()
header(slide, "Next Steps", "Spring 2026")

steps = [
    ("Standardize nuclear mask", "Re-run all models using Cellpose 3 nuclei masks → fair PC comparison"),
    ("3D volume estimation", "Aggregate masks across Z-slices; compute true voxel volumes per condensate"),
    ("CLI packaging", "Replace hardcoded paths in cellpose_pipeline.py with argparse / config file; add docstrings"),
    ("Spring poster / report", "Narrative: problem → Winter pipeline → model survey → Cellpose 3 → 3D volume"),
]

y = Inches(1.5)
for i, (title, desc) in enumerate(steps, 1):
    num_box = slide.shapes.add_shape(1, Inches(0.5), y + Inches(0.02), Inches(0.42), Inches(0.42))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = ACCENT
    num_box.line.fill.background()
    txbox(slide, str(i), Inches(0.5), y, Inches(0.42), Inches(0.45),
          size=17, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txbox(slide, title, Inches(1.05), y, Inches(11.2), Inches(0.38),
          size=17, bold=True, color=WHITE)
    txbox(slide, desc, Inches(1.05), y + Inches(0.36), Inches(11.2), Inches(0.38),
          size=14, color=MUTED)
    y += Inches(1.0)

# ── Slide 7: Reference PC Investigation ──────────────────────────────────────
slide = add_slide()
header(slide, "Reference PC Investigation", "2026-04-23 — found a methodological mismatch")

# Left column: the three PC types
txbox(slide, "Three types of PC", Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.38),
      size=15, bold=True, color=ACCENT)

pc_rows = [
    ("Nuclear PC", "Condensate vs. dilute, both inside the nucleus. What the pipeline computes."),
    ("Cytoplasmic PC", "Condensate vs. dilute, both in the cytoplasm. Not currently computed."),
    ("Background-subtracted PC", "Same formula but a background constant B (outside all cells)\nsubtracted before the ratio. Reference CSVs use this."),
]
y = Inches(1.9)
for title, desc in pc_rows:
    bar = slide.shapes.add_shape(1, Inches(0.5), y + Inches(0.05), Inches(0.06), Inches(0.75))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    txbox(slide, title, Inches(0.68), y, Inches(5.7), Inches(0.32), size=13, bold=True, color=LIGHT)
    txbox(slide, desc, Inches(0.68), y + Inches(0.3), Inches(5.7), Inches(0.55), size=11, color=MUTED)
    y += Inches(0.98)

# Divider
div = slide.shapes.add_shape(1, Inches(6.55), Inches(1.35), Pt(2), Inches(5.8))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC, 0xDD, 0xEE); div.line.fill.background()

# Right column: findings
txbox(slide, "Findings for Sample2_5_1", Inches(6.7), Inches(1.45), Inches(6.1), Inches(0.38),
      size=15, bold=True, color=ACCENT)

findings = [
    ("Dilute density anomaly",
     "Sample2_5_1 nuclear dilute density = 104.09 vs. ~70 for all other\n"
     "Sample2_5 entries. If dilute were ~70, reference PC ≈ 9.4, not 6.32.\n"
     "Possible cause: loose manual mask left bright condensate-edge pixels\n"
     "counted as dilute, deflating the reference PC."),
    ("Background subtraction not in pipeline",
     "Reference CSVs apply background subtraction before computing PC.\n"
     "Pipeline does not. Formula: PC = (cond − B) / (dilute − B)\n"
     "Since dilute ≈ background, subtracting B inflates PC disproportionately.\n"
     "Part of the 3.33 → 6.32 gap is methodological, not just segmentation."),
]
y = Inches(1.9)
for title, desc in findings:
    bar = slide.shapes.add_shape(1, Inches(6.7), y + Inches(0.05), Inches(0.06), Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    txbox(slide, title, Inches(6.88), y, Inches(5.9), Inches(0.32), size=13, bold=True, color=LIGHT)
    txbox(slide, desc, Inches(6.88), y + Inches(0.3), Inches(5.9), Inches(0.95), size=11, color=MUTED)
    y += Inches(1.4)

# Action item
act = slide.shapes.add_shape(1, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.65))
act.fill.solid(); act.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8); act.line.color.rgb = ACCENT
txbox(slide, "Action: confirm background estimation method with lab member → add bg subtraction to pipeline → recompute PC",
      Inches(0.65), Inches(6.6), Inches(12.0), Inches(0.55), size=13, bold=True, color=LIGHT)

out = "spring_checkin_slides.pptx"
prs.save(out)
print(f"Saved: {out}")
